"""
@description: 全格式文档导入 - 支持文本/图片/音频/视频/PPT/Excel/PDF(含嵌入图片)/Word(含嵌入图片)
@dependencies: pathlib, json, src.tools.knowledge_base, src.utils.model_gateway
@last_modified: 2026-03-21
"""
import sys
import io
import json
import base64
import re
import os
from pathlib import Path
from datetime import datetime

# Windows 控制台 UTF-8 编码（仅在终端中执行）
if sys.platform == "win32" and sys.stdout.isatty():
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

sys.path.insert(0, str(Path(__file__).parent.parent))

# 加载 .env 文件
env_file = Path(__file__).parent.parent / ".env"
if env_file.exists():
    for line in env_file.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if line and not line.startswith("#") and "=" in line:
            key, value = line.split("=", 1)
            os.environ.setdefault(key.strip(), value.strip())

from src.tools.knowledge_base import add_knowledge, get_knowledge_stats
from scripts.litellm_gateway import get_model_gateway

INBOX_DIR = Path(__file__).parent.parent / ".ai-state" / "inbox"
PROCESSED_DIR = INBOX_DIR / "processed"

SUPPORTED_TEXT = {".txt", ".md", ".csv", ".json", ".log"}
SUPPORTED_IMAGE = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
SUPPORTED_AUDIO = {".mp3", ".wav", ".ogg", ".m4a", ".flac"}
SUPPORTED_VIDEO = {".mp4", ".mov", ".avi", ".mkv", ".webm"}
SUPPORTED_OFFICE = {".pptx", ".xlsx", ".docx", ".pdf"}
ALL_SUPPORTED = SUPPORTED_TEXT | SUPPORTED_IMAGE | SUPPORTED_AUDIO | SUPPORTED_VIDEO | SUPPORTED_OFFICE


def _read_text_file(filepath: Path) -> str:
    try:
        return filepath.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return filepath.read_text(encoding="gbk", errors="ignore")


def _read_image(filepath: Path) -> str:
    """用 Gemini Vision 理解图片内容"""
    gateway = get_model_gateway()
    image_bytes = filepath.read_bytes()
    ext = filepath.suffix.lower()
    mime_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp"}
    mime = mime_map.get(ext, "image/png")

    result = gateway.call_gemini_vision(
        "gemini_3_pro", image_bytes,
        f"请详细描述这张图片的内容。如果包含文字请全部转写。如果是产品图、设计图、技术图纸，请详细描述产品特征、设计细节和技术参数。文件名：{filepath.name}",
        "", "doc_import_image"
    )
    if result.get("success"):
        return result["response"]
    return f"[图片识别失败: {result.get('error', '')}]"


def _read_audio(filepath: Path) -> str:
    """用 Gemini Audio 理解音频内容"""
    gateway = get_model_gateway()
    audio_bytes = filepath.read_bytes()

    result = gateway.call_gemini_audio(
        "gemini_3_pro", audio_bytes,
        "请转写这段音频的完整内容。如果是中文就输出中文。保留所有有价值的信息。",
        "", "doc_import_audio"
    )
    if result.get("success"):
        return result["response"]
    return f"[音频识别失败: {result.get('error', '')}]"


def _read_video(filepath: Path) -> str:
    """用 Gemini 理解视频内容（提取关键帧+音频）"""
    gateway = get_model_gateway()

    # 视频文件可能很大，限制 20MB
    file_size = filepath.stat().st_size
    if file_size > 20 * 1024 * 1024:
        return f"[视频文件过大: {file_size // 1024 // 1024}MB，上限 20MB]"

    video_bytes = filepath.read_bytes()
    video_b64 = base64.b64encode(video_bytes).decode("utf-8")
    ext = filepath.suffix.lower()
    mime_map = {".mp4": "video/mp4", ".mov": "video/quicktime", ".avi": "video/x-msvideo",
                ".mkv": "video/x-matroska", ".webm": "video/webm"}
    mime = mime_map.get(ext, "video/mp4")

    # 直接调 Gemini API
    import os
    import requests
    api_key = os.environ.get("GEMINI_API_KEY", "")
    if not api_key:
        return "[视频处理失败: 无 GEMINI_API_KEY]"

    payload = {
        "contents": [{"parts": [
            {"inline_data": {"mime_type": mime, "data": video_b64}},
            {"text": f"请详细描述这个视频的内容。包括视觉内容和语音/对话内容。如果涉及产品、技术演示、设计展示，请详细记录。文件名：{filepath.name}"}
        ]}],
        "generationConfig": {"maxOutputTokens": 4096, "temperature": 0.2}
    }

    try:
        url = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent"
        resp = requests.post(f"{url}?key={api_key}", json=payload, timeout=300,
                            headers={"Content-Type": "application/json"}).json()
        if 'candidates' in resp:
            return resp['candidates'][0]['content']['parts'][0]['text']
        return f"[视频处理失败: {str(resp)[:300]}]"
    except Exception as e:
        return f"[视频处理失败: {e}]"


def _read_image_bytes(img_bytes: bytes, label: str = "image") -> str:
    """用 Gemini Vision 理解图片字节"""
    gateway = get_model_gateway()
    result = gateway.call_gemini_vision(
        "gemini_3_pro", img_bytes,
        f"请描述这张图片的内容。如果包含文字请转写。如果是产品图或技术图，详细描述特征。来源：{label}",
        "", "doc_import_embedded_image"
    )
    if result.get("success"):
        return result["response"]
    return ""


def _compress_image_bytes(img_bytes: bytes, max_size: int = 1024) -> bytes:
    """压缩图片字节：限制最大边 max_size 像素，JPEG 质量 85"""
    try:
        from PIL import Image
        from io import BytesIO
        img = Image.open(BytesIO(img_bytes))
        img.thumbnail((max_size, max_size))
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        buf = BytesIO()
        img.save(buf, format="JPEG", quality=85)
        return buf.getvalue()
    except Exception:
        return img_bytes  # 压缩失败就用原图


def _is_garbled(text: str) -> bool:
    """检测文字是否乱码（高比例不可打印字符）"""
    if not text:
        return True
    printable = sum(1 for c in text if c.isprintable() or c in '\n\r\t')
    ratio = printable / len(text)
    return ratio < 0.7


def _read_pdf(filepath: Path) -> str:
    """读取 PDF：提取文字 + 嵌入图片用 Vision 理解"""
    texts = []

    try:
        import fitz  # pymupdf
        doc = fitz.open(str(filepath))

        for page_num, page in enumerate(doc):
            # 提取文字
            page_text = page.get_text().strip()

            # 检测乱码：如果乱码则用整页光栅化 OCR
            if _is_garbled(page_text) or len(page_text) < 50:
                # 整页光栅化，用 Vision OCR
                try:
                    mat = fitz.Matrix(2, 2)  # 2x 放大提高 OCR 精度
                    pix = page.get_pixmap(matrix=mat)
                    img_bytes = pix.tobytes("png")
                    if len(img_bytes) > 5000:
                        desc = _read_image_bytes(img_bytes, f"PDF第{page_num+1}页OCR")
                        if desc and not desc.startswith("["):
                            texts.append(f"[第{page_num+1}页OCR]\n{desc}")
                except Exception:
                    pass
            else:
                texts.append(f"[第{page_num+1}页文字]\n{page_text}")

            # 提取嵌入的大图（非整页渲染的图片）
            images = page.get_images(full=True)
            for img_idx, img_info in enumerate(images):
                try:
                    xref = img_info[0]
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n > 4:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img_bytes = pix.tobytes("png")
                    # 只处理中等大小图片（排除小图标和整页渲染）
                    if 5000 < len(img_bytes) < 500000:
                        desc = _read_image_bytes(img_bytes, f"PDF第{page_num+1}页图片{img_idx+1}")
                        if desc and not desc.startswith("["):
                            texts.append(f"[第{page_num+1}页图片{img_idx+1}]\n{desc}")
                except Exception:
                    continue

            if page_num >= 30:  # 最多处理 30 页
                texts.append(f"[...后续 {len(doc) - 30} 页省略]")
                break
        doc.close()

    except ImportError:
        texts.append("[需要安装 pymupdf: pip install pymupdf --break-system-packages]")
    except Exception as e:
        texts.append(f"[PDF 读取异常: {e}]")

    return "\n\n".join(texts) if texts else f"[PDF 内容为空: {filepath.name}]"


def _read_docx(filepath: Path) -> str:
    """读取 Word：提取文字 + 嵌入图片用 Vision 理解"""
    texts = []

    try:
        from docx import Document

        doc = Document(str(filepath))

        # 提取段落文字
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text.strip())

        # 提取表格
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                table_text.append(" | ".join(row_text))
            if table_text:
                texts.append("[表格]\n" + "\n".join(table_text))

        # 提取嵌入图片
        img_count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    img_bytes = rel.target_part.blob
                    if len(img_bytes) > 1000:
                        img_count += 1
                        desc = _read_image_bytes(img_bytes, f"Word文档图片{img_count}")
                        if desc and not desc.startswith("["):
                            texts.append(f"[嵌入图片{img_count}]\n{desc}")
                except Exception:
                    continue

    except ImportError:
        texts.append("[需要安装 python-docx: pip install python-docx --break-system-packages]")
    except Exception as e:
        texts.append(f"[DOCX 读取异常: {e}]")

    return "\n\n".join(texts) if texts else f"[DOCX 内容为空: {filepath.name}]"


def _read_pptx(filepath: Path) -> str:
    """读取 PPTX：流式逐页处理，多级降级，不跳过任何文件"""
    import gc
    file_size_mb = filepath.stat().st_size / 1024 / 1024

    # 大文件提示
    if file_size_mb > 50:
        print(f"[DocImport] 大文件处理中: {filepath.name} ({file_size_mb:.0f}MB)")

    # === 策略 1: python-pptx 逐页流式处理 ===
    texts = []
    try:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        total_slides = len(prs.slides)
        image_count = 0

        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []

            # 提取文字
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                # 提取表格
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            slide_texts.append(" | ".join(row_text))

            if slide_texts:
                texts.append(f"[第{slide_idx+1}/{total_slides}页]\n" + "\n".join(slide_texts))

            # 提取图片（所有图片都处理，压缩后发 Vision）
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        img_bytes = shape.image.blob
                        if len(img_bytes) > 2000:  # 跳过小图标
                            # 压缩图片
                            compressed = _compress_image_bytes(img_bytes)
                            desc = _read_image_bytes(compressed, f"PPT第{slide_idx+1}页图片")
                            if desc and not desc.startswith("["):
                                texts.append(f"[第{slide_idx+1}页图片]\n{desc}")
                                image_count += 1
                    except Exception:
                        continue

            # 每 10 页释放一次内存
            if (slide_idx + 1) % 10 == 0:
                gc.collect()

        prs = None
        gc.collect()

        if texts:
            print(f"[DocImport] python-pptx 成功: {len(texts)} 段, {image_count} 张图片, {total_slides} 页")
            return "\n\n".join(texts)

    except MemoryError:
        print(f"[DocImport] python-pptx OOM ({file_size_mb:.0f}MB)，降级到 PDF 转换")
        gc.collect()
    except Exception as e:
        print(f"[DocImport] python-pptx 失败: {e}，降级到 PDF 转换")
        gc.collect()

    # === 策略 2: LibreOffice 转 PDF → pymupdf 逐页读 ===
    try:
        import subprocess
        import tempfile

        with tempfile.TemporaryDirectory() as tmp_dir:
            # LibreOffice 转 PDF
            cmd = [
                "soffice", "--headless", "--convert-to", "pdf",
                "--outdir", tmp_dir, str(filepath)
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)

            # 找到生成的 PDF
            pdf_files = list(Path(tmp_dir).glob("*.pdf"))
            if pdf_files:
                pdf_path = pdf_files[0]
                print(f"[DocImport] LibreOffice 转 PDF 成功: {pdf_path.name}")

                # 用 pymupdf 逐页读（内存友好）
                import fitz
                doc = fitz.open(str(pdf_path))
                pdf_texts = []

                for page_num in range(len(doc)):
                    page = doc[page_num]
                    page_text = page.get_text().strip()

                    if page_text and len(page_text) > 20:
                        pdf_texts.append(f"[第{page_num+1}页]\n{page_text}")
                    else:
                        # 文字少的页面用 Vision OCR
                        try:
                            mat = fitz.Matrix(2, 2)
                            pix = page.get_pixmap(matrix=mat)
                            img_bytes = pix.tobytes("png")
                            if len(img_bytes) > 5000:
                                desc = _read_image_bytes(img_bytes, f"PDF第{page_num+1}页")
                                if desc and not desc.startswith("["):
                                    pdf_texts.append(f"[第{page_num+1}页OCR]\n{desc}")
                        except Exception:
                            pass

                    # 每 10 页释放内存
                    if (page_num + 1) % 10 == 0:
                        gc.collect()

                doc.close()
                gc.collect()

                if pdf_texts:
                    print(f"[DocImport] PDF 读取成功: {len(pdf_texts)} 页")
                    return "\n\n".join(pdf_texts)
            else:
                print(f"[DocImport] LibreOffice 未生成 PDF, returncode={result.returncode}")
                if result.stderr:
                    print(f"  stderr: {result.stderr[:300]}")

    except FileNotFoundError:
        print("[DocImport] LibreOffice 未安装，跳过 PDF 降级")
    except subprocess.TimeoutExpired:
        print("[DocImport] LibreOffice 转换超时(300s)")
    except Exception as e:
        print(f"[DocImport] PDF 降级失败: {e}")

    gc.collect()

    # === 策略 3: 直接解压 pptx（它是 zip），读 XML 提取纯文本 ===
    try:
        import zipfile
        import xml.etree.ElementTree as ET

        zip_texts = []
        with zipfile.ZipFile(str(filepath), 'r') as z:
            # 找所有 slide XML
            slide_files = sorted([f for f in z.namelist() if f.startswith("ppt/slides/slide") and f.endswith(".xml")])

            for slide_file in slide_files:
                try:
                    with z.open(slide_file) as sf:
                        tree = ET.parse(sf)
                        root = tree.getroot()

                        # 提取所有文本节点
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        page_texts = []
                        for t_elem in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t'):
                            if t_elem.text and t_elem.text.strip():
                                page_texts.append(t_elem.text.strip())

                        if page_texts:
                            slide_num = slide_file.split("slide")[-1].replace(".xml", "")
                            zip_texts.append(f"[第{slide_num}页]\n" + "\n".join(page_texts))
                except Exception:
                    continue

        gc.collect()

        if zip_texts:
            print(f"[DocImport] ZIP/XML 提取成功: {len(zip_texts)} 页（纯文本，无图片）")
            return "\n\n".join(zip_texts)

    except Exception as e:
        print(f"[DocImport] ZIP/XML 提取失败: {e}")

    # === 全部失败：标记待重试 ===
    pending_dir = Path(__file__).parent.parent / ".ai-state" / "pending_imports"
    pending_dir.mkdir(parents=True, exist_ok=True)

    # 记录失败信息
    import json
    fail_record = {
        "filename": filepath.name,
        "size_mb": round(file_size_mb, 1),
        "failed_at": datetime.now().isoformat(),
        "reason": "所有策略均失败(python-pptx/LibreOffice+PDF/ZIP+XML)",
        "original_path": str(filepath)
    }
    fail_path = pending_dir / f"{filepath.stem}_pending.json"
    fail_path.write_text(json.dumps(fail_record, ensure_ascii=False, indent=2), encoding="utf-8")

    print(f"[DocImport] 全部策略失败，已记录待重试: {fail_path}")
    return f"[处理失败，已记录待后续重试: {filepath.name} ({file_size_mb:.0f}MB)]"


def _read_xlsx(filepath: Path) -> str:
    """读取 Excel 表格"""
    texts = []

    try:
        from openpyxl import load_workbook

        wb = load_workbook(str(filepath), read_only=True, data_only=True)

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(max_row=200, values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                if any(r.strip() for r in row_text):
                    rows.append(" | ".join(row_text))
            if rows:
                texts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()

    except ImportError:
        texts.append("[需要安装 openpyxl: pip install openpyxl --break-system-packages]")
    except Exception as e:
        texts.append(f"[XLSX 读取异常: {e}]")

    return "\n\n".join(texts) if texts else f"[XLSX 内容为空: {filepath.name}]"


def _read_file(filepath: Path) -> str:
    """统一文件读取入口"""
    ext = filepath.suffix.lower()

    if ext in SUPPORTED_TEXT:
        return _read_text_file(filepath)
    elif ext in SUPPORTED_IMAGE:
        return _read_image(filepath)
    elif ext in SUPPORTED_AUDIO:
        return _read_audio(filepath)
    elif ext in SUPPORTED_VIDEO:
        return _read_video(filepath)
    elif ext == ".pdf":
        return _read_pdf(filepath)
    elif ext == ".docx":
        return _read_docx(filepath)
    elif ext == ".pptx":
        return _read_pptx(filepath)
    elif ext == ".xlsx":
        return _read_xlsx(filepath)
    else:
        return f"[不支持的格式: {ext}]"


def _refine_with_llm(content: str, filename: str, is_reread: bool = False) -> dict:
    """用 LLM 提炼文档内容（单条）"""
    if len(content) < 20:
        return {"success": False, "error": "内容过少"}

    gateway = get_model_gateway()

    # 重读时用更深入的 prompt
    if is_reread:
        extra_instruction = (
            "\n注意：这份文档之前已导入过。这次请你以更深的视角重新阅读，关注：\n"
            "1. 之前可能遗漏的细节数据和边缘信息\n"
            "2. 与其他已知知识的新关联（如竞品对比、技术趋势印证）\n"
            "3. 文档中隐含的未说明的假设和风险\n"
            "4. 如果有修改痕迹，标注变化要点\n"
        )
    else:
        extra_instruction = ""

    prompt = (
        f"以下是文件「{filename}」的内容。\n{extra_instruction}\n"
        f"请分析这份文档，提取与「智能骑行头盔产品研发」相关的所有关键信息。\n"
        f"按以下 JSON 格式回复，不要有其他内容：\n"
        f'{{"title": "知识条目标题(20字以内)", '
        f'"domain": "competitors或components或standards或lessons中选一个最合适的", '
        f'"tags": ["标签1", "标签2", "标签3"], '
        f'"summary": "500字以内的结构化摘要，保留关键数据、参数、结论", '
        f'"relevance": "high或medium或low，与智能骑行头盔的相关度"}}\n\n'
        f"如果文档与智能骑行头盔完全无关，relevance 设为 low。\n\n"
        f"文档内容：\n{content[:6000]}"
    )

    result = gateway.call_azure_openai(
        "cpo", prompt,
        "你是研发情报分析专家。只输出 JSON，不输出任何其他内容。",
        "doc_import"
    )

    if not result.get("success"):
        return {"success": False, "error": result.get("error", "LLM 调用失败")}

    response = result["response"].strip()
    response = re.sub(r'^```json\s*', '', response)
    response = re.sub(r'\s*```$', '', response)

    try:
        data = json.loads(response)
        return {"success": True, "data": data}
    except Exception as e:
        return {"success": False, "error": f"JSON 解析失败: {e}"}


def _refine_multi(content: str, filename: str, is_reread: bool = False) -> list:
    """从文档中提取多条知识条目（而非只提取一条摘要）"""
    if len(content) < 50:
        return []

    gateway = get_model_gateway()

    extra = ""
    if is_reread:
        extra = "\n注意：这是重新导入的文档，请关注之前可能遗漏的细节。\n"

    prompt = (
        f"以下是文件「{filename}」的内容。{extra}\n\n"
        f"请将这份文档拆分成多条独立的知识条目。每条知识应该聚焦一个具体的主题/产品/技术点。\n"
        f"例如：一份竞品分析报告应该按每个竞品拆分；一份技术文档应该按技术模块拆分。\n\n"
        f"按以下 JSON 数组格式回复：\n"
        f'[{{"title": "具体标题(含品牌/型号/技术名)", '
        f'"domain": "competitors或components或standards或lessons", '
        f'"tags": ["标签1", "标签2"], '
        f'"summary": "300字以内的结构化摘要，保留关键数据、参数、价格、结论", '
        f'"relevance": "high或medium或low"}}]\n\n'
        f"要求：\n"
        f"1. 尽可能多拆分，每条聚焦一个主题，不要合并\n"
        f"2. 保留所有具体数据（型号、价格、参数、供应商名）\n"
        f"3. 如果文档内容丰富，可以拆出 10-20 条\n"
        f"4. 如果文档很短，拆出 1-3 条也可以\n"
        f"5. domain 只能是 competitors/components/standards/lessons 四选一\n\n"
        f"文档内容（前 8000 字）：\n{content[:8000]}"
    )

    result = gateway.call_azure_openai(
        "cpo", prompt,
        "你是研发情报分析专家。只输出 JSON 数组。不要有其他内容。",
        "doc_import_multi"
    )

    if not result.get("success"):
        return []

    response = result["response"].strip()
    response = re.sub(r'^```json\s*', '', response)
    response = re.sub(r'\s*```$', '', response)

    try:
        items = json.loads(response)
        if isinstance(items, list):
            # 过滤 domain 白名单
            valid_domains = {"competitors", "components", "standards", "lessons"}
            for item in items:
                if item.get("domain") not in valid_domains:
                    item["domain"] = "lessons"
            return items
        elif isinstance(items, dict):
            if items.get("domain") not in {"competitors", "components", "standards", "lessons"}:
                items["domain"] = "lessons"
            return [items]
    except Exception:
        pass
    return []


import threading
_import_lock = threading.Lock()


def scan_and_import(progress_callback=None) -> str:
    """扫描 inbox 目录，处理所有新文件"""
    if not _import_lock.acquire(blocking=False):
        print("[DocImport] 上一次导入仍在进行，跳过")
        return "导入进行中，请稍后"
    try:
        INBOX_DIR.mkdir(parents=True, exist_ok=True)
        PROCESSED_DIR.mkdir(parents=True, exist_ok=True)

        files = [f for f in INBOX_DIR.iterdir() if f.is_file() and f.suffix.lower() in ALL_SUPPORTED]

        if not files:
            return ""

        report_lines = [f"[DocImport] 文档导入 ({datetime.now().strftime('%H:%M')})"]
        report_lines.append(f"[Info] 发现 {len(files)} 个文件\n")
        imported = 0
        skipped = 0

        for f in files:
            short_name = f.name[:30] + "..." if len(f.name) > 30 else f.name

            # 大文件预通知
            file_size_mb = f.stat().st_size / 1024 / 1024
            if file_size_mb > 50 and progress_callback:
                progress_callback(f"[BigFile] 正在处理大文件: {short_name} ({file_size_mb:.0f}MB)，请稍候...")

            if progress_callback:
                progress_callback(f"[Process] {short_name}")

            print(f"[DocImport] 读取: {f.name} ({f.suffix})")
            content = _read_file(f)

            if len(content) < 20 or content.startswith("[不支持") or content.startswith("[需要安装"):
                report_lines.append(f"  [SKIP] {short_name} -- {content[:80]}")
                skipped += 1
                f.rename(PROCESSED_DIR / f.name)
                continue

            print(f"[DocImport] 内容 {len(content)} 字，提炼中...")

            # 检测是否重新导入
            is_reread = (PROCESSED_DIR / f.name).exists()
            if is_reread:
                # 删除旧的 processed 文件，允许新文件进入
                old_processed = PROCESSED_DIR / f.name
                old_processed.unlink(missing_ok=True)
                print(f"[DocImport] 重新导入: {f.name}")

            # 多条目拆分提炼
            items = _refine_multi(content, f.name, is_reread=is_reread)

            if not items:
                report_lines.append(f"  [FAIL] {short_name} -- LLM 提炼失败")
                skipped += 1
                f.rename(PROCESSED_DIR / f.name)
                continue

            # 去重：同标题不重复入库
            seen_titles = set()
            unique_items = []
            for item in items:
                title = item.get("title", "")[:60]
                if title and title not in seen_titles:
                    seen_titles.add(title)
                    unique_items.append(item)
            items = unique_items

            item_count = 0
            for item in items:
                if item.get("relevance") == "low":
                    continue

                add_knowledge(
                    title=item.get("title", f.stem)[:80],
                    domain=item.get("domain", "lessons"),
                    content=item.get("summary", "")[:800],
                    tags=item.get("tags", []),
                    source=f"doc_import:{f.name}",
                    confidence="high" if is_reread or item.get("relevance") == "high" else "medium"
                )
                item_count += 1

            imported += item_count
            report_lines.append(f"  [OK] {short_name} -> {item_count} 条知识")
            f.rename(PROCESSED_DIR / f.name)

            # 每个文件处理完强制释放内存
            import gc
            gc.collect()

        stats = get_knowledge_stats()
        report_lines.append(f"\n[Stats] 知识库现状: {stats}")
        report_lines.append(f"[Summary] 导入: {imported} 条 | 跳过: {skipped} 条")

        report = "\n".join(report_lines)
        print(report)

        # 任务完成提示音
        from src.utils.notifier import notify
        notify("success")

        return report
    finally:
        _import_lock.release()


if __name__ == "__main__":
    report = scan_and_import()
    if not report:
        print("[Info] inbox 为空，无文件需要处理")